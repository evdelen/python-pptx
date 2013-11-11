# encoding: utf-8

"""
Gherkin step implementations for font color features
"""

from __future__ import absolute_import, print_function, unicode_literals

import os

from behave import given, then, when

from pptx import Presentation
from pptx.dml.core import RGBColor
from pptx.enum import MSO_COLOR_TYPE, MSO_THEME_COLOR

from .helpers import saved_pptx_path, test_pptx


font_color_pptx_path = test_pptx('font-color')


# given ===================================================

@given('a font with {color_type} color')
def step_given_font_with_color_type(context, color_type):
    context.textbox_idx = {
        'no':      0,
        'an RGB':  1,
        'a theme': 2
    }[color_type]
    context.prs = Presentation(font_color_pptx_path)
    textbox = context.prs.slides[0].shapes[context.textbox_idx]
    context.font = textbox.textframe.paragraphs[0].runs[0].font


@given('a font with a color brightness setting of {setting}')
def step_font_with_color_brightness(context, setting):
    textbox_idx = {
        'no brightness adjustment': 2,
        '25% darker': 3,
        '40% lighter': 4,
    }[setting]
    context.prs = Presentation(font_color_pptx_path)
    textbox = context.prs.slides[0].shapes[textbox_idx]
    context.font = textbox.textframe.paragraphs[0].runs[0].font


# when ====================================================

@when('I save and reload the presentation')
def step_save_and_reload(context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.prs.save(saved_pptx_path)
    context.prs = Presentation(saved_pptx_path)
    textbox = context.prs.slides[0].shapes[context.textbox_idx]
    context.font = textbox.textframe.paragraphs[0].runs[0].font


@when('I set the font color brightness value')
def step_set_font_color_brightness(context):
    assert False


@when('I set the font {color_type} value')
def step_set_font_color_value(context, color_type):
    if color_type == 'RGB':
        context.font.color.rgb = RGBColor(0x12, 0x34, 0x56)
    elif color_type == 'theme color':
        context.font.color.theme_color = MSO_THEME_COLOR.DARK_1


# then ====================================================

@then('its color value matches its RGB color')
def step_color_value_matches_RGB_color(context):
    assert context.font.color.rgb == RGBColor(255, 102, 0)


@then('its color value matches its theme color')
def step_color_value_matches_theme_color(context):
    assert context.font.color.theme_color == MSO_THEME_COLOR.ACCENT_1


@then("the font's color type is {color_type}")
def step_then_font_color_type_is_value(context, color_type):
    expected_value = {
        'None':        None,
        'RGB':         MSO_COLOR_TYPE.RGB,
        'theme color': MSO_COLOR_TYPE.SCHEME,
    }[color_type]
    assert context.font.color.type == expected_value


@then('its color brightness value is {value}')
def step_color_brightness_value_matches(context, value):
    assert context.font.color.brightness == float(value)


@then("the font's {color_type} value matches the value I set")
def step_color_type_value_matches(context, color_type):
    if color_type == 'RGB':
        assert context.font.color.rgb == RGBColor(0x12, 0x34, 0x56)
    else:
        assert context.font.color.theme_color == MSO_THEME_COLOR.DARK_1


@then("the font's color brightness matches the value I set")
def step_color_brightness_matches(context):
    assert False
